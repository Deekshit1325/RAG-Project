import os
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


def extract_pdf(filepath):
    """Pull text out of a PDF, page by page."""
    pages = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append({
                    "text": text,
                    "page": i + 1,
                    "source": os.path.basename(filepath)
                })
    return pages


def extract_docx(filepath):
    """Word docs don't really have 'pages' so we just grab all paragraphs."""
    doc = Document(filepath)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)

    # also grab text from tables because people put stuff in tables a lot
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                full_text.append(" | ".join(row_text))

    combined = "\n".join(full_text)
    if not combined.strip():
        return []

    return [{
        "text": combined,
        "page": 1,  # no real page numbers in docx
        "source": os.path.basename(filepath)
    }]


def extract_pptx(filepath):
    """Extract text from each slide separately so we can cite slide numbers."""
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
            # also check tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        texts.append(" | ".join(row_data))

        if texts:
            slides.append({
                "text": "\n".join(texts),
                "page": i + 1,  # slide number
                "source": os.path.basename(filepath)
            })
    return slides


def extract_xlsx(filepath):
    """Read each sheet as a separate 'page'. Converts rows to pipe-separated text."""
    wb = load_workbook(filepath, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            # filter out None values and convert everything to string
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows_text.append(" | ".join(cells))

        if rows_text:
            sheets.append({
                "text": "\n".join(rows_text),
                "page": sheet_name,  # using sheet name instead of number
                "source": os.path.basename(filepath)
            })
    return sheets


def extract_txt(filepath):
    """Plain text files — just read and return."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if not text.strip():
        return []
    return [{
        "text": text,
        "page": 1,
        "source": os.path.basename(filepath)
    }]


# maps extension to the right function
EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".txt": extract_txt,
}

SUPPORTED_EXTENSIONS = list(EXTRACTORS.keys())


def load_document(filepath):
    """
    Main entry point — give it a file path, it figures out the type
    and returns a list of dicts like:
    [{"text": "...", "page": 1, "source": "file.pdf"}, ...]
    """
    ext = os.path.splitext(filepath)[1].lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")

    extractor = EXTRACTORS[ext]
    return extractor(filepath)
